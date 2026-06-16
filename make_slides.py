from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

BASE = "/Users/hongseogju/Desktop/2026_1_uni/데이터마이닝/Data Analysis Project"
IMG  = f"{BASE}/notebooks"
OUT  = f"{BASE}/발표_슬라이드.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── 색상 ──────────────────────────────────────────────
C_DARK  = RGBColor(0x1C, 0x28, 0x33)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_NAVY  = RGBColor(0x1A, 0x52, 0x76)
C_BLUE  = RGBColor(0x21, 0x8D, 0xBB)
C_RED   = RGBColor(0xE7, 0x4C, 0x3C)
C_GRAY  = RGBColor(0x5D, 0x6D, 0x7E)
C_GREEN = RGBColor(0x27, 0xAE, 0x60)
C_LGRAY = RGBColor(0xF5, 0xF6, 0xFA)

FONT = 'Apple SD Gothic Neo'

# ── 헬퍼 ──────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txt(slide, text, l, t, w, h,
        size=20, color=C_DARK, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT
    return tb

def img(slide, path, l, t, w=None, h=None):
    if not os.path.exists(path):
        return
    kw = {}
    if w: kw['width']  = Inches(w)
    if h: kw['height'] = Inches(h)
    slide.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def bar(slide, title, bar_color=C_NAVY):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bar_color
    shape.line.fill.background()
    txt(slide, title, 0.3, 0.1, 12.5, 0.9, size=28, color=C_WHITE, bold=True)

def placeholder(slide, label, l, t, w, h):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE0, 0xE7, 0xEF)
    shape.line.color.rgb = C_NAVY
    shape.line.width = Pt(1.5)
    txt(slide, label, l+0.1, t + h/2 - 0.3, w-0.2, 0.6,
        size=16, color=C_NAVY, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════
# Slide 1 — 훅
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)
txt(s, "1,251명이 가입했고", 1.0, 1.8, 11.33, 1.5,
    size=44, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "오늘 활성 사용자는 92명입니다.", 1.0, 3.3, 11.33, 1.5,
    size=44, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "— 81%가 이미 떠났습니다 —", 1.0, 5.0, 11.33, 0.8,
    size=24, color=C_RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# Slide 2 — 앱 소개
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "도서 SNS 앱이란?")
txt(s, "사용자 여정", 0.5, 1.3, 6.0, 0.5, size=16, color=C_GRAY, bold=True)
flow = "가입  →  책 추가  →  아카이브(완독 처리)  →  리뷰 작성\n\n              ↕\n다른 사람 프로필 방문  /  팔로우"
txt(s, flow, 0.5, 1.9, 6.2, 3.5, size=18, color=C_DARK)
placeholder(s, "[ 스크린샷 삽입 ]\n앱 UI 화면", 7.0, 1.3, 5.8, 4.5)

# ══════════════════════════════════════════════════════
# Slide 3 — 데이터 소개
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "분석 데이터")
rows = [
    ("데이터", "규모"),
    ("사용자", "1,077명"),
    ("방문 기록", "5,042건"),
    ("리뷰", "1,666개"),
    ("이탈 피드백", "208건"),
    ("분석 기간", "앱 출시 ~ 2026.05"),
]
for i, (a, b) in enumerate(rows):
    y = 1.3 + i * 0.85
    is_header = i == 0
    color = C_WHITE if is_header else C_DARK
    row_color = C_NAVY if is_header else (RGBColor(0xEB,0xF5,0xFB) if i%2==0 else C_WHITE)
    shape = s.shapes.add_shape(1, Inches(1.5), Inches(y), Inches(4.5), Inches(0.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = row_color
    shape.line.fill.background()
    shape2 = s.shapes.add_shape(1, Inches(6.0), Inches(y), Inches(3.5), Inches(0.8))
    shape2.fill.solid(); shape2.fill.fore_color.rgb = row_color
    shape2.line.fill.background()
    txt(s, a, 1.6, y+0.1, 4.3, 0.6, size=18, color=color, bold=is_header)
    txt(s, b, 6.1, y+0.1, 3.3, 0.6, size=18, color=color, bold=is_header)

# ══════════════════════════════════════════════════════
# Slide 4 — 오늘의 질문
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)
txt(s, "오늘의 질문", 0.5, 0.4, 12.33, 0.8, size=22, color=C_BLUE, bold=True)
questions = ["왜 떠나는가?", "누가 남는가?", "어떻게 막을 수 있는가?"]
y_pos = [1.5, 3.0, 4.5]
for q, y in zip(questions, y_pos):
    txt(s, q, 1.5, y, 10.33, 1.2, size=38, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# Slide 5 — EDA
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "① EDA: 누가 있는가?")
img(s, f"{IMG}/fig_01_active_vs_dormant.png", 0.5, 1.2, w=8.0)
txt(s, "long_dormant 80.9%\n→ 사용자 10명 중 8명은 이미 장기 이탈",
    8.8, 2.5, 4.2, 2.5, size=18, color=C_DARK)
txt(s, "81%", 9.5, 1.3, 3.0, 1.2, size=52, color=C_RED, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# Slide 6 — 코호트
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "② 코호트: 언제 떠나는가?   ✅ H5 지지")
img(s, f"{IMG}/fig_01_ad_cohort_status.png",  0.3, 1.2, w=6.3)
img(s, f"{IMG}/fig_01_cohort_retention.png",  6.8, 1.2, w=6.2)
txt(s, "광고 코호트가 오히려 더 빠르게 이탈\n→ 유입 문제가 아니라 정착 문제",
    0.5, 6.3, 12.33, 0.9, size=17, color=C_NAVY, bold=True)

# ══════════════════════════════════════════════════════
# Slide 7 — 클러스터링
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "③ 클러스터링: 어떤 사람들인가?")
img(s, f"{IMG}/fig_02_persona_radar.png", 0.3, 1.2, w=7.5)
personas = [
    ("소셜 독서가",  "67명  (6%)",   "모든 행동 높음"),
    ("조용한 다수",  "922명 (86%)", "아무것도 안 함 ← 핵심"),
    ("기록형 독자",  "83명  (8%)",  "책·아카이브만"),
]
for i, (name, count, desc) in enumerate(personas):
    y = 1.5 + i * 1.7
    bold = (i == 1)
    c = C_RED if i == 1 else C_DARK
    txt(s, name,  8.0, y,       5.0, 0.5, size=20, color=c, bold=bold)
    txt(s, count, 8.0, y+0.45,  5.0, 0.4, size=16, color=C_GRAY)
    txt(s, desc,  8.0, y+0.85,  5.0, 0.5, size=15, color=c, italic=bold)

# ══════════════════════════════════════════════════════
# Slide 8 — 네트워크
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "④ 네트워크: 연결이 중요한가?   ⚠️ H1 수정")
img(s, f"{IMG}/fig_04_h1_h3.png", 0.3, 1.2, w=7.0)
rows2 = [
    ("지표", "active", "dormant", "차이"),
    ("인디그리 (받은 방문)", "7.21", "7.09", "≈ 없음"),
    ("아웃디그리 (준 방문)", "11.65", "1.32", "8.8배 ★"),
]
for i, (a,b,c,d) in enumerate(rows2):
    y = 1.5 + i*1.1
    is_h = i==0
    bg_c = C_NAVY if is_h else (RGBColor(0xEB,0xF5,0xFB) if i%2 else C_WHITE)
    fc = C_WHITE if is_h else C_DARK
    for xi, (val, wid) in enumerate([(a,3.5),(b,1.5),(c,1.5),(d,1.5)]):
        lx = 7.5 + sum([3.5,1.5,1.5,1.5][:xi])
        sh = s.shapes.add_shape(1,Inches(lx),Inches(y),Inches(wid if xi>0 else 3.5),Inches(1.0))
        sh.fill.solid(); sh.fill.fore_color.rgb = bg_c
        sh.line.fill.background()
        em = C_RED if (val == "8.8배 ★") else fc
        txt(s, val, lx+0.05, y+0.1, wid-0.1, 0.8, size=15, color=em, bold=(is_h or val=="8.8배 ★"))
txt(s, "받는 것보다 주는 행동(능동적 탐색)이 활성을 결정한다",
    7.5, 5.0, 5.5, 0.8, size=16, color=C_NAVY, bold=True)

# ══════════════════════════════════════════════════════
# Slide 9 — 휴면 예측
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "⑤ 휴면 예측: 7일 만에 알 수 있는가?")
img(s, f"{IMG}/fig_03_roc_pr.png",  0.3, 1.2, w=6.3)
img(s, f"{IMG}/fig_03_shap.png",    6.8, 1.2, w=6.2)
txt(s, "AUC 0.648  (목표 0.75 미달)\n→ 이것 자체가 발견: 단발 개입이 아닌 지속 모니터링 필요",
    0.5, 6.2, 12.33, 1.0, size=17, color=C_NAVY, bold=True)

# ══════════════════════════════════════════════════════
# Slide 10 — 브리지
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)
txt(s, "7일로는 못 잡는다.", 1.0, 2.0, 11.33, 1.5,
    size=42, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "그렇다면 무엇이 살아남게 하는가?", 1.0, 3.7, 11.33, 1.5,
    size=36, color=C_BLUE, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# Slide 11 — 생존 분석 KM
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "⑥ 생존 분석: 무엇이 오래 남게 하는가?")
img(s, f"{IMG}/fig_06_km_curves.png", 0.3, 1.2, w=9.0)
verdicts = [
    ("H2  책 수 多", "❌ 기각  p=0.751"),
    ("H3  팔로잉 有", "✅ 지지  p<0.001"),
    ("H4  아카이브 有", "✅ 지지  p<0.001"),
]
for i, (h, v) in enumerate(verdicts):
    y = 1.8 + i * 1.5
    c = C_GREEN if "✅" in v else C_RED
    txt(s, h, 9.6, y,      3.5, 0.5, size=15, color=C_DARK, bold=True)
    txt(s, v, 9.6, y+0.45, 3.5, 0.6, size=17, color=c, bold=True)

# ══════════════════════════════════════════════════════
# Slide 12 — Cox
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "⑦ Cox 모델: 행동별 이탈 위험 감소율")
img(s, f"{IMG}/fig_06_cox.png", 0.3, 1.2, w=6.5)
cox_rows = [
    ("행동",         "HR",    "위험 감소"),
    ("리뷰 작성",    "0.715", "28.5% ★"),
    ("아카이브",     "0.744", "25.6%"),
    ("프로필 방문",  "0.799", "20.1%"),
    ("책 추가",      "0.821", "17.9%"),
]
for i, (a,b,c) in enumerate(cox_rows):
    y = 1.5 + i * 1.0
    is_h = i==0
    bg_c = C_NAVY if is_h else (RGBColor(0xEB,0xF5,0xFB) if i%2 else C_WHITE)
    fc = C_WHITE if is_h else C_DARK
    for xi, (val, wid) in enumerate([(a,3.5),(b,1.5),(c,2.0)]):
        lx = 7.0 + sum([3.5,1.5,2.0][:xi])
        sh = s.shapes.add_shape(1,Inches(lx),Inches(y),Inches(wid),Inches(0.9))
        sh.fill.solid(); sh.fill.fore_color.rgb = bg_c
        sh.line.fill.background()
        em = C_RED if "★" in val else fc
        txt(s, val, lx+0.05, y+0.1, wid-0.1, 0.7, size=16, color=em,
            bold=(is_h or "★" in val))

# ══════════════════════════════════════════════════════
# Slide 13 — 이탈 피드백
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "⑧ 이탈 피드백: 사용자는 뭐라고 했는가?")
img(s, f"{IMG}/fig_05_feedback.png", 0.3, 1.2, w=7.5)
txt(s, "1위  사용빈도가 낮아서  40%",
    8.0, 1.8, 5.0, 0.7, size=17, color=C_GRAY)
txt(s, "→ 증상이지 원인이 아님",
    8.0, 2.4, 5.0, 0.5, size=14, color=C_GRAY, italic=True)
txt(s, "2위  기록 기능 부적합  19%",
    8.0, 3.2, 5.0, 0.7, size=17, color=C_RED, bold=True)
txt(s, "→ 앱의 핵심 가치(기록)를\n    제대로 경험시키지 못함",
    8.0, 3.85, 5.0, 1.0, size=14, color=C_RED)
txt(s, "정량 분석: 기록 행동이 리텐션 핵심\n이탈자 목소리: 기록 경험이 부족했다",
    7.8, 5.4, 5.3, 1.5, size=16, color=C_NAVY, bold=True)

# ══════════════════════════════════════════════════════
# Slide 14 — 가설 종합
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "가설 검증 종합")
hypo = [
    ("가설", "결과", "판정"),
    ("H1: 받은 방문 多 → 활성",      "아웃디그리가 8.8배 더 중요", "⚠️ 수정"),
    ("H2: 책 수 多 → 생존",          "p = 0.751 (n.s.)",           "❌ 기각"),
    ("H3: 팔로잉 有 → 생존",         "p < 0.001",                  "✅ 지지"),
    ("H4: 아카이브 → 생존",          "HR = 0.744, p < 0.001",      "✅ 지지"),
    ("H5: 광고 코호트 리텐션 낮음",   "코호트 비교 확인",            "✅ 지지"),
]
for i, (a,b,c) in enumerate(hypo):
    y = 1.3 + i * 0.97
    is_h = i==0
    bg_c = C_NAVY if is_h else (C_WHITE if i%2 else RGBColor(0xEB,0xF5,0xFB))
    fc = C_WHITE if is_h else C_DARK
    vc = (C_GREEN if "✅" in c else C_RED if "❌" in c else RGBColor(0xE6,0x7E,0x22)) if not is_h else C_WHITE
    for xi, (val, wid) in enumerate([(a,6.0),(b,4.5),(c,2.0)]):
        lx = 0.5 + sum([6.0,4.5,2.0][:xi])
        sh = s.shapes.add_shape(1,Inches(lx),Inches(y),Inches(wid),Inches(0.9))
        sh.fill.solid(); sh.fill.fore_color.rgb = bg_c
        sh.line.fill.background()
        col = vc if xi==2 else fc
        txt(s, val, lx+0.1, y+0.1, wid-0.15, 0.7, size=16, color=col, bold=(is_h or xi==2))

# ══════════════════════════════════════════════════════
# Slide 15 — Aha Moment
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)
txt(s, "6개 분석이 하나를 가리킨다", 0.5, 0.5, 12.33, 0.8,
    size=22, color=C_BLUE, align=PP_ALIGN.CENTER)
txt(s, "Aha Moment", 0.5, 1.4, 12.33, 1.2,
    size=52, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, "첫 주에 아카이브 또는 리뷰를\n한 번이라도 경험한 사용자",
    0.5, 2.8, 12.33, 1.8,
    size=34, color=RGBColor(0xF7,0xDC,0x6F), bold=True, align=PP_ALIGN.CENTER)
bullets = [
    "SHAP: 아카이브 중요도 1위",
    "Cox: 리뷰 HR 0.715  /  아카이브 HR 0.744",
    "클러스터링: 86%가 이 경험 없이 이탈",
    "이탈 피드백: 19% '기록 기능 부적합'",
]
for i, b in enumerate(bullets):
    txt(s, f"• {b}", 2.5, 4.8 + i*0.5, 8.33, 0.5,
        size=17, color=C_GRAY)

# ══════════════════════════════════════════════════════
# Slide 16 — 액션 아이템
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_LGRAY)
bar(s, "액션 아이템")
actions = [
    ("1", "온보딩에서 첫 아카이브 필수 유도",     "SHAP 1위  /  Cox HR 0.744"),
    ("2", "가입 7일 내 리뷰 작성 푸시 알림",       "Cox 최강 보호 요인 HR 0.715"),
    ("3", "탐색 피드 강화 (능동적 방문 유도)",      "아웃디그리 8.8배 효과"),
]
for i, (num, act, basis) in enumerate(actions):
    y = 1.5 + i * 1.8
    circle = s.shapes.add_shape(9, Inches(0.5), Inches(y), Inches(0.9), Inches(0.9))
    circle.fill.solid(); circle.fill.fore_color.rgb = C_NAVY
    circle.line.fill.background()
    txt(s, num, 0.5, y+0.1, 0.9, 0.7, size=26, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, act,   1.6, y,       11.0, 0.6, size=22, color=C_DARK, bold=True)
    txt(s, basis, 1.6, y+0.6,   11.0, 0.5, size=16, color=C_GRAY)

# ══════════════════════════════════════════════════════
# Slide 17 — 한계 & 마무리
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_DARK)
txt(s, "한계 및 향후 과제", 0.5, 0.3, 12.33, 0.8,
    size=26, color=C_BLUE, bold=True)
limits = [
    ("AUC 0.648",           "7일 예측 한계 → 30일 누적 피처로 재모델링 필요"),
    ("이탈 텍스트 98% 결측", "질적 분석 불가 → 탈퇴 시 서술 피드백 수집 권고"),
    ("랜딩페이지 데이터 없음","A/B 테스트 데이터 구축 필요"),
]
for i, (title, desc) in enumerate(limits):
    y = 1.3 + i * 1.5
    txt(s, f"• {title}", 0.8, y,      11.5, 0.6, size=20, color=C_WHITE, bold=True)
    txt(s, f"  {desc}",  0.8, y+0.55, 11.5, 0.6, size=17, color=C_GRAY)

txt(s, "감사합니다", 1.0, 5.5, 11.33, 1.5,
    size=44, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

# ── 저장 ──────────────────────────────────────────────
prs.save(OUT)
print(f"저장 완료: {OUT}")
print(f"총 슬라이드: {len(prs.slides)}장")
